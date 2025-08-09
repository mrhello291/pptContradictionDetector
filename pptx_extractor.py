"""
PowerPoint content extraction module.
"""
import io
import re
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.table import Table
from PIL import Image
import base64

from models import SlideContent

logger = logging.getLogger(__name__)

class PPTXExtractor:
    """Extracts content from PowerPoint presentations."""
    
    def __init__(self):
        self.numerical_patterns = [
            r'\$?[\d,]+\.?\d*[kKmMbBtT]?',  # Currency and large numbers
            r'\d+\.?\d*%',  # Percentages
            r'\d{4}[-/]\d{1,2}[-/]\d{1,2}',  # Dates
            r'\d+\.?\d*',  # General numbers
        ]
    
    def extract_presentation(self, file_path: Path) -> List[SlideContent]:
        """Extract content from all slides in a PowerPoint presentation."""
        logger.info(f"Extracting content from {file_path}")
        
        try:
            presentation = Presentation(file_path)
            slides = []
            
            for slide_idx, slide in enumerate(presentation.slides, 1):
                slide_content = self._extract_slide_content(slide, slide_idx)
                slides.append(slide_content)
                logger.debug(f"Extracted content from slide {slide_idx}")
            
            logger.info(f"Successfully extracted content from {len(slides)} slides")
            return slides
            
        except Exception as e:
            logger.error(f"Error extracting presentation: {e}")
            raise
    
    def _extract_slide_content(self, slide, slide_number: int) -> SlideContent:
        """Extract content from a single slide."""
        title = self._extract_title(slide)
        text_content = self._extract_text_content(slide)
        numerical_data = self._extract_numerical_data(text_content)
        images_text = self._extract_images_text(slide)
        raw_content = self._get_raw_content(slide)
        
        return SlideContent(
            slide_number=slide_number,
            title=title,
            text_content=text_content,
            numerical_data=numerical_data,
            images_text=images_text,
            raw_content=raw_content
        )
    
    def _extract_title(self, slide) -> Optional[str]:
        """Extract the title from a slide."""
        try:
            if hasattr(slide, 'shapes') and slide.shapes.title:
                return slide.shapes.title.text.strip()
        except AttributeError:
            pass
        return None
    
    def _extract_text_content(self, slide) -> List[str]:
        """Extract all text content from a slide."""
        text_content = []
        
        for shape in slide.shapes:
            text = self._extract_text_from_shape(shape)
            if text:
                text_content.append(text)
        
        return text_content
    
    def _extract_text_from_shape(self, shape: BaseShape) -> str:
        """Extract text from a shape, including tables."""
        text_parts = []
        
        # Handle text frames
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = ''.join(run.text for run in paragraph.runs)
                if paragraph_text.strip():
                    text_parts.append(paragraph_text.strip())
        
        # Handle tables (GraphicFrame with table)
        elif hasattr(shape, 'table') and shape.table is not None:
            table_text = self._extract_table_text(shape.table)
            if table_text:
                text_parts.append(table_text)
        
        # Handle grouped shapes
        elif hasattr(shape, 'shapes'):
            for subshape in shape.shapes:
                subtext = self._extract_text_from_shape(subshape)
                if subtext:
                    text_parts.append(subtext)
        
        return ' '.join(text_parts)
    
    def _extract_table_text(self, table) -> str:
        """Extract text from a table."""
        table_text = []
        
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = ' '.join(paragraph.text for paragraph in cell.text_frame.paragraphs)
                row_text.append(cell_text.strip())
            if any(row_text):
                table_text.append(' | '.join(row_text))
        
        return '\n'.join(table_text)
    
    def _extract_numerical_data(self, text_content: List[str]) -> List[Dict[str, Any]]:
        """Extract numerical data from text content."""
        numerical_data = []
        
        for text in text_content:
            for pattern in self.numerical_patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                for match in matches:
                    numerical_data.append({
                        'value': match,
                        'context': text,
                        'type': self._classify_numerical_data(match)
                    })
        
        return numerical_data
    
    def _classify_numerical_data(self, value: str) -> str:
        """Classify the type of numerical data."""
        if '%' in value:
            return 'percentage'
        elif '$' in value or any(suffix in value.lower() for suffix in ['k', 'm', 'b', 't']):
            return 'currency'
        elif re.match(r'\d{4}[-/]\d{1,2}[-/]\d{1,2}', value):
            return 'date'
        else:
            return 'number'
    
    def _extract_images_text(self, slide) -> List[str]:
        """Extract text from images in a slide (placeholder for OCR)."""
        # This is a placeholder - in a real implementation, you would:
        # 1. Extract images from the slide
        # 2. Use OCR (like Tesseract) to extract text
        # 3. Return the extracted text
        
        images_text = []
        
        for shape in slide.shapes:
            if isinstance(shape, Picture):
                # Placeholder: In real implementation, perform OCR here
                # For now, we'll just note that an image exists
                images_text.append("[IMAGE_DETECTED - OCR would extract text here]")
        
        return images_text
    
    def _get_raw_content(self, slide) -> str:
        """Get raw content of the slide for comprehensive analysis."""
        all_text = []
        
        # Extract all text content
        for shape in slide.shapes:
            text = self._extract_text_from_shape(shape)
            if text:
                all_text.append(text)
        
        return '\n'.join(all_text)
    
    def save_slide_images(self, file_path: Path, output_dir: Path) -> List[Path]:
        """Save slide images for image-based analysis (if needed)."""
        # This is a placeholder for extracting slide images
        # In a real implementation, you might use python-pptx or convert to images
        logger.info("Image extraction not implemented in this version")
        return [] 
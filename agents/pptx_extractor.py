"""
PowerPoint content extraction module.
"""
import io
import re
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple, ByteString
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.table import Table
from PIL import Image
import base64

from utils.models import SlideContent

logger = logging.getLogger(__name__)

class PPTXExtractor:
    """Extracts content from PowerPoint presentations."""
    
    def __init__(self):
        self.numerical_patterns = [
            r'\$?[\d,]+\.?\d*[kKmMbBtT]?',  # Currency and large numbers
            r'\d+\.?\d*%',  # Percentages
            r'\d{4}[-/]\d{1,2}[-/]\d{1,2}',  # Dates
            r'\d+\.?\d*',  # General numbers
        ]
    
    def extract_presentation(self, file_path: Path) -> List[SlideContent]:
        """Extract content from all slides in a PowerPoint presentation."""
        logger.info(f"Extracting content from {file_path}")
        
        try:
            presentation = Presentation(file_path)
            slides = []
            
            # New Step: Render all slides to images once
            slides_as_images = self._render_slides_to_images(file_path)
            
            for slide_idx, slide in enumerate(presentation.slides, 1):
                # Pass the full slide image data to the next method
                slide_image_data = slides_as_images[slide_idx - 1]
                
                slide_content = self._extract_slide_content(slide, slide_idx, slide_image_data)
                slides.append(slide_content)
                logger.debug(f"Extracted content from slide {slide_idx}")
            
            logger.info(f"Successfully extracted content from {len(slides)} slides")
            return slides
            
        except Exception as e:
            logger.error(f"Error extracting presentation: {e}")
            raise
    
    def _extract_slide_content(self, slide, slide_number: int, slide_image_data: ByteString) -> SlideContent:
        """Extract content from a single slide."""
        title = self._extract_title(slide)
        text_content = self._extract_text_content(slide)
        numerical_data = self._extract_numerical_data(text_content)
        # images_data = self._extract_images_data(slide)
        raw_content = self._get_raw_content(slide)
        
        # New Step: Collect all image data (full slide + individual images)
        images_data = [slide_image_data] if slide_image_data else []
        images_data.extend(self._extract_individual_images_data(slide))
        
        return SlideContent(
            slide_number=slide_number,
            title=title,
            text_content=text_content,
            numerical_data=numerical_data,
            images_data=images_data,
            raw_content=raw_content
        )
    
    def _extract_title(self, slide) -> Optional[str]:
        """Extract the title from a slide."""
        try:
            if hasattr(slide, 'shapes') and slide.shapes.title:
                return slide.shapes.title.text.strip()
        except AttributeError:
            pass
        return None
    
    def _extract_text_content(self, slide) -> List[str]:
        """Extract all text content from a slide."""
        text_content = []
        
        for shape in slide.shapes:
            text = self._extract_text_from_shape(shape)
            if text:
                text_content.append(text)
        
        return text_content
    
    def _extract_text_from_shape(self, shape) -> str:
        """Extract text from a shape, handling all shape types safely."""
        text_parts = []
        
        try:
            # Handle text frames (most common text container)
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    if paragraph_text.strip():
                        text_parts.append(paragraph_text.strip())
            
            # Handle tables (GraphicFrame with table) - safely check
            if hasattr(shape, 'table'):
                try:
                    if shape.table is not None:
                        table_text = self._extract_table_text(shape.table)
                        if table_text:
                            text_parts.append(table_text)
                except (ValueError, AttributeError):
                    # Not a table shape, skip
                    pass
            
            # Handle charts (GraphicFrame with chart)
            if hasattr(shape, 'chart'):
                try:
                    if shape.chart is not None:
                        # Extract chart title and data labels if available
                        chart_text = self._extract_chart_text(shape.chart)
                        if chart_text:
                            text_parts.append(chart_text)
                except (ValueError, AttributeError):
                    # Not a chart shape, skip
                    pass
            
            # Handle grouped shapes
            if hasattr(shape, 'shapes'):
                try:
                    for subshape in shape.shapes:
                        subtext = self._extract_text_from_shape(subshape)
                        if subtext:
                            text_parts.append(subtext)
                except (AttributeError, TypeError):
                    # Not a group shape, skip
                    pass
            
            # Handle text directly on shape (for some shape types)
            if hasattr(shape, 'text') and not text_parts:
                try:
                    if shape.text and shape.text.strip():
                        text_parts.append(shape.text.strip())
                except (AttributeError, TypeError):
                    pass
                    
        except Exception as e:
            # Log the error but don't crash
            logger.debug(f"Error extracting text from shape {type(shape)}: {e}")
        
        return ' '.join(text_parts)
    
    def _extract_table_text(self, table) -> str:
        """Extract text from a table."""
        table_text = []
        
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = ' '.join(paragraph.text for paragraph in cell.text_frame.paragraphs)
                row_text.append(cell_text.strip())
            if any(row_text):
                table_text.append(' | '.join(row_text))
        
        return '\n'.join(table_text)
    
    def _extract_numerical_data(self, text_content: List[str]) -> List[Dict[str, Any]]:
        """Extract numerical data from text content."""
        numerical_data = []
        
        for text in text_content:
            for pattern in self.numerical_patterns:
                matches = re.findall(pattern, text, re.IGNORECASE)
                for match in matches:
                    numerical_data.append({
                        'value': match,
                        'context': text,
                        'type': self._classify_numerical_data(match)
                    })
        
        return numerical_data
    
    def _classify_numerical_data(self, value: str) -> str:
        """Classify the type of numerical data."""
        if '%' in value:
            return 'percentage'
        elif '$' in value or any(suffix in value.lower() for suffix in ['k', 'm', 'b', 't']):
            return 'currency'
        elif re.match(r'\d{4}[-/]\d{1,2}[-/]\d{1,2}', value):
            return 'date'
        else:
            return 'number'
    
    def _extract_individual_images_data(self, slide) -> List[ByteString]:
        """Extract raw image data from all individual picture shapes in a slide."""
        individual_images = []
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                individual_images.append(shape.image.blob)
        return individual_images
    
    def _render_slides_to_images(self, file_path: Path) -> List[ByteString]:
        """
        Helper function to render each slide of the presentation as an image.
        This is a placeholder for your chosen implementation (e.g., using LibreOffice, a library, etc.).
        """
        logger.warning("Rendering slides to images is not implemented. Please add your chosen method here.")
        # You would need to add your implementation here.
        # Example using a hypothetical library:
        # from pptx_image import convert_pptx_to_to_images
        # return convert_pptx_to_to_images(file_path)
        
        # Returning a placeholder for now to prevent crashes
        return [b""] * len(Presentation(file_path).slides)
    
    def _get_raw_content(self, slide) -> str:
        """Get raw content of the slide for comprehensive analysis."""
        all_text = []
        
        # Extract all text content
        for shape in slide.shapes:
            text = self._extract_text_from_shape(shape)
            if text:
                all_text.append(text)
        
        return '\n'.join(all_text)
    
    def save_slide_images(self, file_path: Path, output_dir: Path) -> List[Path]:
        """Save slide images for image-based analysis (if needed)."""
        # This is a placeholder for extracting slide images
        # In a real implementation, you might use python-pptx or convert to images
        logger.info("Image extraction not implemented in this version")
        return [] 
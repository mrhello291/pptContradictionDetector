#!/usr/bin/env python3
"""
Create a sample PowerPoint presentation with intentional inconsistencies
for testing the contradiction detector.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

def create_sample_presentation():
    """Create a sample presentation with various types of inconsistencies."""
    
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Company Q3 2024 Financial Report"
    subtitle.text = "Quarterly Business Review\nPrepared by Finance Team\nOctober 2024"
    
    # Slide 2: Revenue Overview
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Q3 Revenue Overview"
    
    content = slide.placeholders[1]
    content.text = """Key Financial Metrics:
• Total Q3 Revenue: $2.4M
• Growth Rate: 15% YoY
• Market Share: 8.5%
• Customer Acquisition: 240 new customers
• Customer Retention: 94%"""
    
    # Slide 3: Market Analysis
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Market Analysis"
    
    content = slide.placeholders[1]
    content.text = """Market Conditions:
• The market is highly competitive with 12 major players
• Our market share increased to 8.5% 
• Customer demand is strong
• Competition is intensifying with new entrants
• Market size: $45B globally"""
    
    # Slide 4: Quarterly Breakdown (INCONSISTENCY: Different Q3 revenue)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Quarterly Revenue Breakdown"
    
    content = slide.placeholders[1]
    content.text = """2024 Quarterly Performance:
• Q1 Revenue: $2.1M
• Q2 Revenue: $2.3M  
• Q3 Revenue: $2.7M (Contradicts slide 2!)
• Q4 Forecast: $2.8M

Total YTD: $7.1M"""
    
    # Slide 5: Customer Metrics (INCONSISTENCY: Different customer numbers)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Customer Metrics Deep Dive"
    
    content = slide.placeholders[1]
    content.text = """Customer Growth Analysis:
• New customers acquired in Q3: 180 (Contradicts slide 2!)
• Customer retention rate: 92% (Different from slide 2!)
• Average customer value: $4,200
• Customer lifetime value: $18,500
• Churn rate: 8% (Should be 6% if retention is 94%)"""
    
    # Slide 6: Market Competition (INCONSISTENCY: Contradicts slide 3)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Competitive Landscape"
    
    content = slide.placeholders[1]
    content.text = """Competitive Environment:
• Market has limited competition with only 3 major competitors (Contradicts slide 3!)
• We dominate with 8.5% market share
• Barriers to entry are high
• Market is consolidating rapidly
• Low competitive pressure"""
    
    # Slide 7: Percentage Issues
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Revenue Distribution by Product"
    
    content = slide.placeholders[1]
    content.text = """Product Line Revenue Share:
• Product A: 35%
• Product B: 28%
• Product C: 22%
• Product D: 18%
• Other: 5%

Total: 108% (Should be 100%!)"""
    
    # Slide 8: Timeline Issues
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Key Milestones & Timeline"
    
    content = slide.placeholders[1]
    content.text = """Important Dates:
• Company founded: March 2020
• First major client: January 2019 (Before company was founded!)
• Product launch: June 2021
• Series A funding: December 2020
• IPO planned: Q2 2025
• Market expansion: Q1 2024 (Already happened but listed as future)"""
    
    # Slide 9: Summary (More inconsistencies)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Executive Summary"
    
    content = slide.placeholders[1]
    content.text = """Q3 2024 Highlights:
• Strong revenue growth: Q3 revenue of $2.9M (Third different number!)
• Excellent customer metrics: 220 new customers
• Market leadership in low-competition environment
• Successfully navigated highly competitive market
• 97% customer retention (Fourth different number!)
• Ready for continued growth"""
    
    return prs

def main():
    """Create and save the sample presentation."""
    print("🎯 Creating sample PowerPoint with intentional inconsistencies...")
    
    try:
        # Create presentation
        prs = create_sample_presentation()
        
        # Save presentation
        output_file = Path("sample_inconsistent_presentation.pptx")
        prs.save(str(output_file))
        
        print(f"✅ Sample presentation created: {output_file}")
        print()
        print("📋 Inconsistencies included:")
        print("• Slide 2 vs 4: Q3 revenue ($2.4M vs $2.7M)")  
        print("• Slide 2 vs 5: New customers (240 vs 180)")
        print("• Slide 2 vs 5: Retention rate (94% vs 92%)")
        print("• Slide 3 vs 6: Market competition (12 players vs 3)")
        print("• Slide 5: Retention vs churn rate (92% retention but 8% churn)")
        print("• Slide 7: Percentages sum to 108% instead of 100%")
        print("• Slide 8: Timeline issues (client before company founding)")
        print("• Slide 9: Third Q3 revenue number ($2.9M)")
        print("• Multiple contradictory claims about market conditions")
        print()
        print("🔍 Test the detector with:")
        print(f"python ppt_contradiction_detector.py {output_file}")
        
    except Exception as e:
        print(f"❌ Error creating sample presentation: {e}")
        print("Make sure you have python-pptx installed: pip install python-pptx")

if __name__ == "__main__":
    main() 